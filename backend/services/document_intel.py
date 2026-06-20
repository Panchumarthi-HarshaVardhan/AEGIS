# ============================================================
# JARVIS Guardian AI — Document Intelligence Service
# Text extraction from PDF/DOCX and injection auditing
# ============================================================

import os
from typing import Dict, Any
from groq import Groq

class DocumentIntelAnalyzer:
    def __init__(self, api_key: str):
        if not api_key:
            raise ValueError("Groq API Key is required for Document Intelligence")
        self.client = Groq(api_key=api_key)
        self.model = "llama-3.3-70b-versatile"

    async def analyze_document(self, file_path: str) -> Dict[str, Any]:
        """Extracts text, runs prompt injection scanner, and summarizes document contents."""
        if not os.path.exists(file_path):
            return {"success": False, "error": "Document file not found"}

        file_name = os.path.basename(file_path)
        ext = file_name.split('.')[-1].lower()
        extracted_text = ""

        try:
            # 1. Text Extraction based on file extension
            if ext == 'pdf':
                extracted_text = self._extract_pdf_text(file_path)
            elif ext == 'docx':
                extracted_text = self._extract_docx_text(file_path)
            elif ext == 'pptx':
                extracted_text = self._extract_pptx_text(file_path)
            elif ext in ['txt', 'md', 'json', 'csv', 'py', 'js', 'ts', 'css', 'html']:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    extracted_text = f.read()
            else:
                return {"success": False, "error": f"Unsupported file extension: {ext}"}

            if not extracted_text.strip():
                return {"success": False, "error": "No readable text content extracted from document."}

            # Limit text size sent to LLM
            content_sample = extracted_text[:8000]

            # 2. Run Security Audits (Prompt Injection checking)
            injection_detected, injection_reason = self._scan_for_injections(content_sample)
            
            # 3. Generate Summary & QA Analysis via Groq
            summary = await self._generate_summary(content_sample)

            return {
                "success": True,
                "file_name": file_name,
                "file_size_bytes": os.path.getsize(file_path),
                "text_length": len(extracted_text),
                "summary": summary,
                "security": {
                    "approved": not injection_detected,
                    "reason": injection_reason if injection_detected else "No malicious prompt instructions found."
                }
            }

        except Exception as e:
            return {"success": False, "error": f"Document analysis failed: {str(e)}"}

    def _extract_pdf_text(self, file_path: str) -> str:
        """Extracts text from PDF file."""
        from pypdf import PdfReader
        text = []
        reader = PdfReader(file_path)
        for page in reader.pages[:15]: # Limit to first 15 pages
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
        return "\n".join(text)

    def _extract_docx_text(self, file_path: str) -> str:
        """Extracts text from Microsoft Word DOCX file."""
        import docx
        doc = docx.Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])

    def _extract_pptx_text(self, file_path: str) -> str:
        """Extracts text from PowerPoint PPTX slides."""
        from pptx import Presentation
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
        return "\n".join(text)

    def _scan_for_injections(self, text: str) -> tuple[bool, str]:
        """Checks text for hidden system prompts, jailbreaks, or command override phrases."""
        patterns = [
            ("ignore previous instructions", "Prompt Injection (Ignore Instructions)"),
            ("system: override", "Prompt Injection (Override context)"),
            ("you are no longer jarvis", "Jailbreak attempt (Persona Hijack)"),
            ("instead of summarizing", "Instruction Hijack attempt")
        ]
        
        lower_text = text.lower()
        for pattern, desc in patterns:
            if pattern in lower_text:
                return True, f"Security alert: Blocked document due to synthetic payload match: '{desc}'"
                
        return False, ""

    async def _generate_summary(self, text: str) -> str:
        """Sends extracted text to Llama 3.3 for structured summarization."""
        prompt = (
            "You are JARVIS Guardian AI's document intelligence reader. Analyze the extracted text below. "
            "Generate a professional, structured document report with:\n"
            "- A 2-sentence executive summary.\n"
            "- 4 primary bullet points listing key findings or facts.\n"
            "- A final security warning note if the text contains sensitive keys, passwords, or risky requests.\n\n"
            f"Extracted Text:\n{text}"
        )
        
        completion = self.client.chat.completions.create(
            model=self.model,
            temperature=0.2,
            messages=[{"role": "user", "content": prompt}]
        )
        
        return completion.choices[0].message.content.strip()
